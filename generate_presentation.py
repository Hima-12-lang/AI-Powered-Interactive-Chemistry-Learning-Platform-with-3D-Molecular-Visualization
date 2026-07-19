from pptx import Presentation

prs = Presentation()
slides = [
    {
        "title": "Chemistry Structure Studio",
        "content": [
            "A web app for exploring chemical structures, tracking learning progress, and user accounts.",
            "Built with React, Vite, Express, PostgreSQL, and JWT authentication."
        ]
    },
    {
        "title": "Key Features",
        "content": [
            "Interactive molecule viewers and structure comparisons.",
            "Signup/login with secure backend authentication.",
            "Progress storage in PostgreSQL for personalized learning.",
            "Responsive UI with modular React components."
        ]
    },
    {
        "title": "Frontend Architecture",
        "content": [
            "React + Vite application in src/ with component-based UI.",
            "Main pages: molecule viewer, learning lab, comparison modal, chat assistant.",
            "Auth form and profile banner integrated with backend API calls."
        ]
    },
    {
        "title": "Backend & Database",
        "content": [
            "Node.js + Express backend serving auth and progress APIs.",
            "PostgreSQL stores user accounts and progress data.",
            "Secure password hashing with bcrypt and JWT tokens."
        ]
    },
    {
        "title": "Deployment & Usage",
        "content": [
            "Run frontend with npm and backend with Node/Express.",
            "Use .env for DB credentials and JWT secret.",
            "App supports signup, login, progress load/save, and logout."
        ]
    },
    {
        "title": "Project Summary",
        "content": [
            "A submission-ready chemistry learning tool with auth and persistence.",
            "Designed to demonstrate full-stack web development skills.",
            "Includes documentation and backend integration for real data storage."
        ]
    }
]

for slide_info in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_info["title"]
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = slide_info["content"][0]
    for line in slide_info["content"][1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0

prs.save("Chemistry-Structure-Studio-Presentation.pptx")
print("Presentation saved: Chemistry-Structure-Studio-Presentation.pptx")
